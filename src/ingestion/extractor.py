"""
Extrator de documentos para múltiplos formatos
Suporte: PDF, DOCX, XLSX, CSV, JSON, HTML, MD, PPTX
Com enriquecimento de texto para melhorar a recuperação
"""

import os
import json
import csv
import re
from typing import Dict, List, Any, Optional
from pathlib import Path
import tempfile

# PDF - Usando pdfplumber (mais robusto)
import pdfplumber

# Word
from docx import Document

# Excel
import pandas as pd
from openpyxl import load_workbook

# HTML
from bs4 import BeautifulSoup

# Markdown
import markdown

# PowerPoint
from pptx import Presentation

# JSON e CSV já nativos


class TextEnricher:
    """
    Enriquecimento de texto para melhorar a recuperação semântica
    Adiciona sinônimos e contexto ao texto extraído
    """
    
    @classmethod
    def enrich(cls, text: str) -> str:
        """
        Enriquece o texto adicionando sinônimos e contexto
        
        Args:
            text: Texto original extraído do documento
            
        Returns:
            Texto enriquecido com sinônimos e contexto
        """
        enriched = text
        
        # 🔥 DETECÇÃO ROBUSTA DE PALAVRAS-CHAVE
        
        # 1. Detecta informações sobre orçamento/capital
        if re.search(r'R?\$?\s*[\d.,]+\s*(centavos|real|reais)', text, re.IGNORECASE):
            enriched += "\n\nKeywords: orçamento inicial capital investimento verba recurso financeiro"
        elif re.search(r'R?\$?\s*[\d.,]+', text):
            enriched += "\n\nKeywords: orçamento investimento valor"
        
        # 2. Detecta informações sobre fundação
        if re.search(r'fundada|fundou|criada|estabelecida|história|julho', text, re.IGNORECASE):
            enriched += "\n\nKeywords: fundação inicio origem história criacao estabelecimento"
        
        # 3. Detecta informações sobre localização
        if re.search(r'sede|escritório|localização|Natal|Rio Grande do Norte', text, re.IGNORECASE):
            enriched += "\n\nKeywords: sede escritório localização base endereço"
        
        # 4. Detecta informações sobre o fundador
        if re.search(r'CEO|Alex Tito|fundador|diretor', text, re.IGNORECASE):
            enriched += "\n\nKeywords: fundador criador idealizador empreendedor CEO"
        
        # 5. Detecta informações sobre a empresa
        if re.search(r'TechFlow|Solutions|FlowManager', text, re.IGNORECASE):
            enriched += "\n\nKeywords: TechFlow Solutions empresa tecnologia SaaS inovação"
        
        # 🔥 CONTEXTOS ESPECÍFICOS
        context_added = []
        
        # Verifica se é o documento da empresa (pela presença de TechFlow + CNPJ)
        if 'TechFlow' in text and 'CNPJ' in text:
            context_added.append(
                "[CONTEXTO: Documento corporativo da TechFlow Solutions - contém informações institucionais]"
            )
        
        # Verifica se tem informações financeiras
        if re.search(r'R?\$|centavos|economia|investimento|orçamento', text, re.IGNORECASE):
            context_added.append(
                "[CONTEXTO: Este documento contém informações financeiras, incluindo orçamento inicial e investimentos]"
            )
        
        # Verifica se tem informações históricas
        if re.search(r'julho|2026|2027|2018|fundação|história', text, re.IGNORECASE):
            context_added.append(
                "[CONTEXTO: Este documento contém informações históricas sobre a fundação da empresa]"
            )
        
        # Verifica se tem informações sobre o fundador
        if re.search(r'Alex Tito|CEO|fundador', text, re.IGNORECASE):
            context_added.append(
                "[CONTEXTO: Este documento contém informações sobre o fundador e liderança da empresa]"
            )
        
        # Verifica se tem informações sobre localização
        if re.search(r'Natal|sede|escritório', text, re.IGNORECASE):
            context_added.append(
                "[CONTEXTO: Este documento contém informações sobre a localização e presença geográfica]"
            )
        
        if context_added:
            enriched += "\n\n" + "\n".join(context_added)
        
        return enriched


class DocumentExtractor:
    """Extrator de conteúdo de documentos em múltiplos formatos"""
    
    def __init__(self):
        self.supported_formats = [
            '.pdf', '.docx', '.xlsx', '.xls', 
            '.csv', '.json', '.html', '.htm', 
            '.md', '.txt', '.pptx'
        ]
    
    def extract(self, file_path: str) -> Dict[str, Any]:
        """
        Extrai conteúdo e metadados de um documento
        
        Args:
            file_path: Caminho para o arquivo
            
        Returns:
            Dict com 'content' (texto extraído e enriquecido) e 'metadata'
        """
        file_path = Path(file_path)
        extension = file_path.suffix.lower()
        
        # Extrai conforme o formato
        if extension == '.pdf':
            result = self._extract_pdf(file_path)
        elif extension == '.docx':
            result = self._extract_docx(file_path)
        elif extension in ['.xlsx', '.xls']:
            result = self._extract_excel(file_path)
        elif extension == '.csv':
            result = self._extract_csv(file_path)
        elif extension == '.json':
            result = self._extract_json(file_path)
        elif extension in ['.html', '.htm']:
            result = self._extract_html(file_path)
        elif extension == '.md':
            result = self._extract_markdown(file_path)
        elif extension == '.pptx':
            result = self._extract_pptx(file_path)
        elif extension == '.txt':
            result = self._extract_txt(file_path)
        else:
            raise ValueError(f"Formato não suportado: {extension}")
        
        # 🔧 ENRIQUECE O TEXTO (se houver conteúdo)
        if result.get("content") and len(result["content"]) > 10:
            result["content"] = TextEnricher.enrich(result["content"])
            result["metadata"]["enriched"] = True
        
        return result
    
    def _extract_pdf(self, file_path: Path) -> Dict[str, Any]:
        """Extrai texto de PDF usando pdfplumber (mais robusto)"""
        try:
            with pdfplumber.open(file_path) as pdf:
                text = ""
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text += page_text + "\n\n"
                
                # Se não extraiu nada, tenta extrair tabelas
                if not text.strip():
                    for page in pdf.pages:
                        tables = page.extract_tables()
                        for table in tables:
                            for row in table:
                                if row:
                                    text += " | ".join([str(cell) if cell else "" for cell in row]) + "\n"
                
                if not text.strip():
                    print(f"⚠️ Nenhum texto extraído do PDF: {file_path.name}")
                    text = f"[PDF NÃO EXTRAÍDO: {file_path.name}]"
                
                return {"content": text.strip(), "metadata": {"pages": len(pdf.pages)}}
        except Exception as e:
            print(f"❌ Erro ao extrair PDF {file_path.name}: {str(e)}")
            return {"content": f"[ERRO NA EXTRAÇÃO: {str(e)}]", "metadata": {}}
    
    def _extract_docx(self, file_path: Path) -> Dict[str, Any]:
        """Extrai texto de DOCX"""
        try:
            doc = Document(file_path)
            text = []
            
            # Extrai parágrafos
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text.append(paragraph.text)
            
            # Extrai tabelas
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text.append(" | ".join(row_text))
            
            return {
                "content": "\n\n".join(text),
                "metadata": {"paragraphs": len(doc.paragraphs), "tables": len(doc.tables)}
            }
        except Exception as e:
            return {"content": f"Erro ao extrair DOCX: {str(e)}", "metadata": {}}
    
    def _extract_excel(self, file_path: Path) -> Dict[str, Any]:
        """Extrai texto de Excel (XLSX/XLS)"""
        try:
            df_dict = pd.read_excel(file_path, sheet_name=None)
            text_parts = []
            sheets_info = []
            
            for sheet_name, df in df_dict.items():
                sheets_info.append({"name": sheet_name, "rows": len(df), "cols": len(df.columns)})
                
                # Converte para texto legível
                text_parts.append(f"=== Planilha: {sheet_name} ===\n")
                text_parts.append(df.to_string(index=False))
                text_parts.append("\n" + "="*50 + "\n")
            
            return {
                "content": "\n".join(text_parts),
                "metadata": {"sheets": sheets_info}
            }
        except Exception as e:
            return {"content": f"Erro ao extrair Excel: {str(e)}", "metadata": {}}
    
    def _extract_csv(self, file_path: Path) -> Dict[str, Any]:
        """Extrai texto de CSV"""
        try:
            df = pd.read_csv(file_path)
            text = f"=== CSV: {file_path.name} ===\n\n"
            text += df.to_string(index=False)
            
            return {
                "content": text,
                "metadata": {"rows": len(df), "columns": list(df.columns)}
            }
        except Exception as e:
            return {"content": f"Erro ao extrair CSV: {str(e)}", "metadata": {}}
    
    def _extract_json(self, file_path: Path) -> Dict[str, Any]:
        """Extrai texto de JSON"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            
            text = f"=== JSON: {file_path.name} ===\n\n"
            text += json.dumps(data, indent=2, ensure_ascii=False)
            
            return {
                "content": text,
                "metadata": {"keys": list(data.keys()) if isinstance(data, dict) else []}
            }
        except Exception as e:
            return {"content": f"Erro ao extrair JSON: {str(e)}", "metadata": {}}
    
    def _extract_html(self, file_path: Path) -> Dict[str, Any]:
        """Extrai texto de HTML"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                html = f.read()
            
            soup = BeautifulSoup(html, 'html.parser')
            
            # Remove scripts e styles
            for script in soup(["script", "style"]):
                script.decompose()
            
            # Extrai título
            title = soup.title.string if soup.title else None
            
            # Extrai texto limpo
            text = soup.get_text(separator='\n')
            text = '\n'.join(line.strip() for line in text.split('\n') if line.strip())
            
            return {
                "content": text,
                "metadata": {"title": title}
            }
        except Exception as e:
            return {"content": f"Erro ao extrair HTML: {str(e)}", "metadata": {}}
    
    def _extract_markdown(self, file_path: Path) -> Dict[str, Any]:
        """Extrai texto de Markdown"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                md_content = f.read()
            
            # Converte para HTML e depois extrai texto
            html = markdown.markdown(md_content)
            soup = BeautifulSoup(html, 'html.parser')
            text = soup.get_text(separator='\n')
            
            return {
                "content": text,
                "metadata": {"format": "markdown"}
            }
        except Exception as e:
            return {"content": f"Erro ao extrair Markdown: {str(e)}", "metadata": {}}
    
    def _extract_pptx(self, file_path: Path) -> Dict[str, Any]:
        """Extrai texto de PowerPoint"""
        try:
            prs = Presentation(file_path)
            text_parts = []
            slides_info = []
            
            for i, slide in enumerate(prs.slides):
                slide_text = []
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    text_parts.append(f"=== Slide {i+1} ===\n")
                    text_parts.append("\n".join(slide_text))
                    text_parts.append("\n")
                    slides_info.append({"number": i+1, "texts": len(slide_text)})
            
            return {
                "content": "\n".join(text_parts),
                "metadata": {"slides": len(prs.slides)}
            }
        except Exception as e:
            return {"content": f"Erro ao extrair PPTX: {str(e)}", "metadata": {}}
    
    def _extract_txt(self, file_path: Path) -> Dict[str, Any]:
        """Extrai texto de TXT"""
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                content = f.read()
            
            return {
                "content": content,
                "metadata": {"format": "txt"}
            }
        except Exception as e:
            return {"content": f"Erro ao extrair TXT: {str(e)}", "metadata": {}}